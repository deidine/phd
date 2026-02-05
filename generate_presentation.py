"""
PhD Thesis Proposal Presentation
Moving Target Defense & Formal Authorization Verification
for Securing Distributed Microservices — NO Machine Learning
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x1A, 0x5C, 0x9E)
TEAL   = RGBColor(0x00, 0xBF, 0xA5)
RED    = RGBColor(0xD8, 0x35, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xD0, 0xEC, 0xF8)
GREY   = RGBColor(0x99, 0xAA, 0xBB)
PURPLE = RGBColor(0x6A, 0x0D, 0xAD)

def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def box(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def t(slide, text, l, top, w, h, size=14, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def bl(slide, items, l, top, w, h, size=12, title=None, tc=TEAL, bc=LIGHT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    if title:
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = title; r.font.size = Pt(size+2)
        r.font.bold = True; r.font.color.rgb = tc; first = False
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]; first = False
        r = p.add_run(); r.text = f"▸  {item}"
        r.font.size = Pt(size); r.font.color.rgb = bc

# SLIDE 1 — TITLE
def s1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, DARK)
    box(sl, 0, 0, 0.4, 7.5, BLUE)
    box(sl, 0.4, 0.22, 9.6, 0.07, TEAL)
    t(sl, "PhD THESIS PROPOSAL  ·  CYBERSECURITY & DISTRIBUTED SYSTEMS",
      0.65, 0.42, 9, 0.45, size=12, color=GREY)
    t(sl, "Moving Target Defense\n& Formal Authorization Verification\nfor Securing Distributed Microservices",
      0.65, 0.9, 9.1, 2.6, size=28, bold=True, color=WHITE)
    t(sl, "No Machine Learning  ·  Algorithmic Security  ·  Formal Methods",
      0.65, 3.55, 9, 0.5, size=15, color=TEAL)
    box(sl, 0.65, 4.35, 5.5, 0.06, RED)
    t(sl, "Deidine Cheigeur  ·  June 2026", 0.65, 4.5, 6, 0.45, size=14, color=GREY)
    t(sl, "Field: Cybersecurity + Distributed Systems",
      0.65, 4.95, 6, 0.4, size=13, color=GREY)

# SLIDE 2 — THE PROBLEM
def s2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, DARK)
    box(sl, 0, 0, 10, 1.1, RED)
    t(sl, "THE PROBLEM", 0.4, 0.3, 9, 0.55, size=26, bold=True)
    t(sl, "Cloud microservices are distributed and dynamic — existing security tools are static and centralised.",
      0.4, 1.2, 9.2, 0.5, size=15, color=LIGHT)
    attacks = [
        ("Reconnaissance", BLUE, "Attacker scans static\nservice IPs & ports\nto map the system"),
        ("Lateral\nMovement",  PURPLE, "Compromised service\npivots to others using\nstolen credentials"),
        ("Privilege\nEscalation", RED, "Attacker gains\nunauthorised permissions\nthrough policy flaws"),
        ("DDoS", RGBColor(0x22,0x7A,0x44), "Flooding one service\ncauses cascade failure\nacross the mesh"),
    ]
    for i, (title, color, body) in enumerate(attacks):
        x = 0.4 + i * 2.4
        box(sl, x, 1.85, 2.15, 0.65, color)
        t(sl, title, x+0.05, 1.9, 2.05, 0.6, size=13, bold=True, align=PP_ALIGN.CENTER)
        t(sl, body, x+0.1, 2.6, 2.0, 1.3, size=12, color=LIGHT)
    box(sl, 0.4, 4.15, 9.2, 0.07, TEAL)
    bl(sl, ["Static endpoints make reconnaissance trivial — attacker maps the system once",
            "Authorization policies are never formally checked for logical flaws",
            "All current detection tools require ML — complex, opaque, hard to deploy in constrained environments"],
       0.4, 4.3, 9.2, 2.2, size=13,
       title="Why current defences fail:", tc=TEAL)

# SLIDE 3 — THESIS OVERVIEW
def s3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, DARK)
    box(sl, 0, 0, 10, 1.1, BLUE)
    t(sl, "THESIS OVERVIEW  —  3 Contributions, 0 Machine Learning",
      0.4, 0.3, 9.2, 0.55, size=23, bold=True)
    cards = [
        ("C1 — Moving Target Defense", BLUE,
         "Kubernetes controller that\nperiodically rotates service\nendpoints, IPs & API paths.\nAttacker maps → system changes."),
        ("C2 — Statistical Detection", TEAL,
         "Shannon entropy + CUSUM\ncontrol charts on authorization\nlogs. Pure math. No ML.\nDetects DDoS & privilege probing."),
        ("C3 — Formal Verification", PURPLE,
         "TLA+ model of Zanzibar/Keto\nauthorization graph.\nProves: no privilege escalation,\nno lateral movement possible."),
    ]
    for i, (title, color, body) in enumerate(cards):
        x = 0.4 + i * 3.2
        box(sl, x, 1.25, 3.0, 0.65, color)
        t(sl, title, x+0.1, 1.3, 2.8, 0.6, size=14, bold=True)
        t(sl, body, x+0.1, 2.05, 2.85, 2.5, size=13, color=LIGHT)
    box(sl, 0.4, 4.75, 9.2, 0.07, RED)
    t(sl, "Authorization layer: Ory Keto (open-source Google Zanzibar) — github.com/ory/keto",
      0.4, 4.9, 9.2, 0.45, size=13, color=RGBColor(0xFF,0xCC,0x55))
    t(sl, "Research question: Can MTD + statistical detection + formally-verified authorization "
          "prevent attacks in distributed microservices without any machine learning?",
      0.4, 5.45, 9.2, 0.9, size=13, color=LIGHT)

# SLIDE 4 — THE RESEARCH GAP
def s4(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, DARK)
    box(sl, 0, 0, 10, 1.1, TEAL)
    t(sl, "THE RESEARCH GAP  —  Why This Has Not Been Done Yet",
      0.4, 0.28, 9.2, 0.6, size=22, bold=True, color=DARK)
    quote = ('"Application of Moving Target Defense to cloud-native, '
             'container-orchestrated environments remains an open research direction. '
             'No existing work integrates MTD with authorization systems '
             'or provides formal security guarantees."')
    box(sl, 0.4, 1.2, 9.2, 1.25, BLUE)
    t(sl, quote, 0.55, 1.3, 8.9, 1.1, size=13, color=WHITE)
    t(sl, "— Sengupta et al. (2020), IEEE Communications Surveys & Tutorials",
      0.55, 2.4, 8.9, 0.38, size=11, color=GREY)
    gaps = [
        ("Gap 1", "No MTD system integrated with\nZanzibar/Keto authorization"),
        ("Gap 2", "Zanzibar (USENIX 2019) has\nnever been formally verified"),
        ("Gap 3", "No statistical (non-ML) IDS\nfor microservices auth layer"),
        ("Gap 4", "No system combines all three:\nMTD + stats + formal proof"),
    ]
    for i, (g, desc) in enumerate(gaps):
        x = 0.4 + (i % 2) * 4.85
        y = 3.1 + (i // 2) * 1.65
        box(sl, x, y, 0.9, 1.35, RED)
        t(sl, g, x+0.05, y+0.35, 0.8, 0.65, size=12, bold=True, align=PP_ALIGN.CENTER)
        t(sl, desc, x+1.05, y+0.2, 3.65, 1.0, size=13, color=LIGHT)

# SLIDE 5 — METHODOLOGY
def s5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, DARK)
    box(sl, 0, 0, 10, 1.1, BLUE)
    t(sl, "METHODOLOGY  —  No Machine Learning", 0.4, 0.3, 9, 0.55, size=26, bold=True)
    steps = [
        ("MTD Engine", BLUE,
         ["Kubernetes Python controller", "Rotates ClusterIPs + ports every 60s",
          "Services found via Keto tuples", "Metric: attacker mean-time-to-exploit"]),
        ("Statistical\nDetector", TEAL,
         ["Shannon entropy on request flows", "CUSUM on per-service request rate",
          "Alert: entropy drop > 40%", "No training data required"]),
        ("TLA+ Formal\nVerification", PURPLE,
         ["Model Keto tuple graph in TLA+", "Invariant 1: no privilege escalation",
          "Invariant 2: no lateral movement", "TLC model checker: proof or bug"]),
    ]
    for i, (title, color, items) in enumerate(steps):
        x = 0.35 + i * 3.2
        box(sl, x, 1.25, 3.0, 0.75, color)
        t(sl, title, x+0.1, 1.3, 2.8, 0.7, size=15, bold=True, align=PP_ALIGN.CENTER)
        if i < 2:
            t(sl, "→", x+3.0, 1.6, 0.2, 0.4, size=18, color=GREY, align=PP_ALIGN.CENTER)
        bl(sl, items, x, 2.1, 3.1, 3.5, size=12, bc=LIGHT)
    box(sl, 0.35, 5.85, 9.3, 0.07, RED)
    t(sl, "Evaluation: CIC-IDS2017 dataset  +  live Kubernetes testbed  +  TLC model checker output",
      0.35, 6.0, 9.3, 0.45, size=13, color=GREY)

# SLIDE 6 — MINI PROJECT
def s6(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, DARK)
    box(sl, 0, 0, 10, 1.1, TEAL)
    t(sl, "MINI PROJECT  —  phd-mtd-zanzibar-security",
      0.4, 0.28, 9.2, 0.6, size=22, bold=True, color=DARK)
    bl(sl,
       ["mtd_controller/     — Kubernetes service endpoint rotation (Python)",
        "stat_detector/      — Entropy + CUSUM anomaly detector",
        "tla_specs/          — TLA+ formal models of Keto authorization",
        "keto_integration/   — Ory Keto policy definitions (OPL)",
        "evaluation/         — CIC-IDS2017 pipeline + testbed attack scripts",
        "tests/              — pytest integration tests"],
       0.4, 1.25, 5.7, 4.6,
       title="Repository Structure", tc=TEAL)
    bl(sl,
       ["Python 3.11", "Kubernetes client (k8s Python lib)",
        "Ory Keto (authorization)", "numpy + scipy (entropy, CUSUM)",
        "TLA+ Toolbox / VS Code", "minikube + kubectl",
        "GitHub Actions CI"],
       6.3, 1.25, 3.4, 4.6,
       title="Tech Stack (No ML)", tc=TEAL)
    box(sl, 0.4, 6.1, 9.2, 0.95, BLUE)
    t(sl,
      "Goal: MTD reduces attacker mean-time-to-exploit by > 80%  ·  "
      "Statistical detector: F1 > 92%  ·  TLA+ proves 2 security invariants",
      0.55, 6.15, 8.9, 0.8, size=12, color=WHITE)

# SLIDE 7 — REFERENCES
def s7(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, DARK)
    box(sl, 0, 0, 10, 1.1, BLUE)
    t(sl, "KEY REFERENCES", 0.4, 0.3, 9, 0.55, size=26, bold=True)
    refs = [
        ("[BOOK]", BLUE,
         "Anderson, R. (2020). Security Engineering, 3rd ed. Wiley. (Free PDF)\n"
         "→ Foundation: access control, formal methods, distributed systems security."),
        ("[A1]", TEAL,
         "Pang et al. (2019). Zanzibar: Google's Consistent, Global Authorization System. USENIX ATC.\n"
         "→ The paper behind Ory Keto. Your authorization platform. Never formally verified — your gap."),
        ("[A2]", PURPLE,
         "Sengupta et al. (2020). A Survey of Moving Target Defenses for Network Security.\n"
         "IEEE Comm. Surveys & Tutorials. → Maps the MTD field; confirms cloud-native MTD is open."),
        ("[A3]", RED,
         "Newcombe et al. (2015). How Amazon Web Services Uses Formal Methods. CACM 58(4).\n"
         "→ Proves TLA+ finds real bugs in distributed systems. Justifies your Contribution 3."),
    ]
    for i, (tag, color, body) in enumerate(refs):
        y = 1.3 + i * 1.5
        box(sl, 0.4, y, 0.8, 1.2, color)
        t(sl, tag, 0.42, y+0.28, 0.75, 0.65, size=12, bold=True, align=PP_ALIGN.CENTER)
        t(sl, body, 1.35, y+0.08, 8.3, 1.2, size=12, color=LIGHT)

# SLIDE 8 — TIMELINE
def s8(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, DARK)
    box(sl, 0, 0, 10, 1.1, BLUE)
    t(sl, "TIMELINE  —  6-Month Catchup + Year 2 Execution",
      0.4, 0.3, 9.2, 0.55, size=22, bold=True)
    phases = [
        ("M1–2\nFoundations",  BLUE,
         ["Read 25 papers", "Literature review", "Thesis Ch.1 & 2"]),
        ("M3\nPrototype",      TEAL,
         ["Build Kubernetes testbed", "MTD controller v1", "CIC-IDS2017 pipeline"]),
        ("M4\nMethodology",    PURPLE,
         ["Thesis Ch. 3", "Statistical detector", "TLA+ model v1"]),
        ("M5–6\nResults",      RED,
         ["Thesis Ch. 4 & 5", "Full evaluation", "Submit conference paper"]),
    ]
    for i, (phase, color, items) in enumerate(phases):
        x = 0.3 + i * 2.4
        box(sl, x, 1.25, 2.15, 0.75, color)
        t(sl, phase, x+0.05, 1.3, 2.05, 0.7, size=13, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            t(sl, "→", x+2.15, 1.55, 0.25, 0.4, size=18, color=GREY, align=PP_ALIGN.CENTER)
        bl(sl, items, x, 2.1, 2.2, 2.8, size=12, bc=LIGHT)
    box(sl, 0.3, 5.1, 9.4, 0.07, TEAL)
    t(sl, "Target journals: IEEE Trans. Dependable & Secure Computing  ·  Computers & Security",
      0.3, 5.25, 9.4, 0.45, size=13, color=GREY)
    box(sl, 0.3, 5.85, 9.4, 1.35, RGBColor(0x0A, 0x2A, 0x45))
    t(sl, "YOUR ELEVATOR PITCH — Memorise this:",
      0.45, 5.9, 9.0, 0.4, size=12, bold=True, color=TEAL)
    t(sl,
      '"My thesis secures distributed microservices using Moving Target Defense '
      '(rotate endpoints in Kubernetes), statistical detection (entropy + CUSUM, no ML), '
      'and TLA+ formal verification of Zanzibar/Keto authorization policies — '
      'proving no privilege escalation or lateral movement is possible."',
      0.45, 6.3, 9.0, 0.85, size=11, color=WHITE)

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    for fn in [s1, s2, s3, s4, s5, s6, s7, s8]:
        fn(prs)
    out = "/Users/deidinecheigeur/Desktop/projects/phd/PhD_Thesis_Proposal_MTD_FormalAuth.pptx"
    prs.save(out)
    print(f"Saved → {out}")

if __name__ == "__main__":
    main()
